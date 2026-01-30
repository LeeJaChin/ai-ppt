import logging
import os
import sys
from pathlib import Path
from typing import Optional
import subprocess
import shutil
import traceback

# 设置库路径以解决 macOS 上的 GDI+ 库问题
if sys.platform == "darwin":  # macOS
    lib_path = "/opt/homebrew/lib"
    if os.path.exists(lib_path):
        os.environ.setdefault("DYLD_LIBRARY_PATH", lib_path + ":" + os.environ.get("DYLD_LIBRARY_PATH", ""))
    else:
        # 对于其他 Homebrew 安装位置的备选方案
        alt_lib_paths = [
            "/usr/local/lib",
            "/opt/homebrew/lib"
        ]
        for path in alt_lib_paths:
            if os.path.exists(path):
                os.environ.setdefault("DYLD_LIBRARY_PATH", path + ":" + os.environ.get("DYLD_LIBRARY_PATH", ""))
                break

try:
    import aspose.slides as slides
except ImportError as e:
    slides = None
    logging.warning(f"Failed to import aspose.slides: {e}")

try:
    import aspose.words as words
except ImportError as e:
    words = None
    logging.warning(f"Failed to import aspose.words: {e}")

# 尝试导入其他可能的转换库
try:
    from pdf2docx import Converter as PdfToDocxConverter
except ImportError:
    PdfToDocxConverter = None

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    Presentation = None
    Inches = None

logger = logging.getLogger(__name__)

class FileConverter:
    """文件转换器类，支持多种文档格式转换"""
    
    def __init__(self):
        self.supported_formats = {
            'ppt_to_pdf': ['.ppt', '.pptx'],
            'word_to_pdf': ['.doc', '.docx'],
            'pdf_to_word': ['.pdf'],
            'pdf_to_ppt': ['.pdf']
        }
    
    def ppt_to_pdf(self, input_path: str, output_path: str) -> bool:
        """
        使用 Aspose.Slides 将 PPT/PPTX 转换为 PDF
        
        Args:
            input_path: 输入 PPT/PPTX 文件路径
            output_path: 输出 PDF 文件路径
            
        Returns:
            bool: 转换是否成功
        """
        logger.info(f"开始转换 PPT 文件: {input_path} -> {output_path}")
        
        if slides is None:
            logger.error("Aspose.Slides not available")
            return False
            
        try:
            # 加载演示文稿
            pres = slides.Presentation(input_path)
            
            # 保存为 PDF
            pres.save(output_path, slides.export.SaveFormat.PDF)
            
            logger.info(f"PPT 转 PDF 成功: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"PPT 转 PDF 失败: {str(e)}")
            # 如果 aspose 失败，尝试使用 unoconv 作为备选方案
            return self._fallback_ppt_to_pdf_with_unoconv(input_path, output_path)
    
    def word_to_pdf(self, input_path: str, output_path: str) -> bool:
        """
        使用 Aspose.Words 将 Word 文档转换为 PDF
        
        Args:
            input_path: 输入 Word 文件路径
            output_path: 输出 PDF 文件路径
            
        Returns:
            bool: 转换是否成功
        """
        logger.info(f"开始转换 Word 文件: {input_path} -> {output_path}")
        
        if words is None:
            logger.error("Aspose.Words not available")
            return False
            
        try:
            # 加载文档
            doc = words.Document(input_path)
            
            # 保存为 PDF
            doc.save(output_path, words.SaveFormat.PDF)
            
            logger.info(f"Word 转 PDF 成功: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"Word 转 PDF 失败: {str(e)}")
            return False
    
    def pdf_to_word(self, input_path: str, output_path: str) -> bool:
        """
        将 PDF 转换为 Word 文档
        
        Args:
            input_path: 输入 PDF 文件路径
            output_path: 输出 Word 文件路径
            
        Returns:
            bool: 转换是否成功
        """
        logger.info(f"开始转换 PDF 文件: {input_path} -> {output_path}")
        
        if PdfToDocxConverter is None:
            logger.error("pdf2docx not available")
            return False
        
        try:
            # 使用 pdf2docx 进行转换
            cv = PdfToDocxConverter(input_path)
            cv.convert(output_path)
            cv.close()
            
            logger.info(f"PDF 转 Word 成功: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"PDF 转 Word 失败: {str(e)}")
            return False
    
    def pdf_to_ppt(self, input_path: str, output_path: str) -> bool:
        """
        将 PDF 转换为 PPT
        
        Args:
            input_path: 输入 PDF 文件路径
            output_path: 输出 PPT 文件路径
            
        Returns:
            bool: 转换是否成功
        """
        logger.info(f"开始转换 PDF 文件: {input_path} -> {output_path}")
        
        if PdfToDocxConverter is None or Presentation is None:
            logger.error("Required libraries for PDF to PPT conversion not available")
            return False
        
        try:
            # 首先将 PDF 转换为 Word，然后将 Word 内容导入 PPT
            temp_docx = input_path.replace('.pdf', '_temp.docx')
            
            # PDF to Word
            cv = PdfToDocxConverter(input_path)
            cv.convert(temp_docx)
            cv.close()
            
            # Word to PPT
            prs = Presentation()
            doc = words.Document(temp_docx)
            
            # 添加每一段文本到幻灯片
            slide = None
            for para in doc.paragraphs:
                if para.text.strip():
                    if slide is None:
                        slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局
                        title = slide.shapes.title
                        title.text = "Converted Slide"
                        
                        content = slide.shapes.placeholders[1]
                        content.text = para.text
                    else:
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        title = slide.shapes.title
                        title.text = f"Slide {len(prs.slides)}"
                        
                        content = slide.shapes.placeholders[1]
                        content.text = para.text
            
            prs.save(output_path)
            
            # 清理临时文件
            if os.path.exists(temp_docx):
                os.remove(temp_docx)
                
            logger.info(f"PDF 转 PPT 成功: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"PDF 转 PPT 失败: {str(e)}")
            return False

def _fallback_ppt_to_pdf_with_unoconv(self, input_path: str, output_path: str) -> bool:
    """
    使用 unoconv 或 LibreOffice (soffice) 作为备选方案进行 PPT 到 PDF 的转换
    先尝试 unoconv，如果不可用或失败再尝试 soffice（libreoffice）。
    """
    logger.info(f"使用 unoconv/soffice 作为备选方案进行 PPT 转 PDF: {input_path} -> {output_path}")
    try:
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        # 1) 尝试 unoconv
        unoconv_path = shutil.which("unoconv")
        if unoconv_path:
            cmd = [unoconv_path, "-f", "pdf", "-o", output_path, input_path]
            logger.debug(f"调用 unoconv: {' '.join(cmd)}")
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
            logger.debug(f"unoconv returncode={result.returncode} stdout={result.stdout} stderr={result.stderr}")
            if result.returncode == 0 and os.path.exists(output_path):
                logger.info(f"Fallback (unoconv) PPT to PDF 成功: {output_path}")
                return True
            else:
                logger.warning(f"unoconv 转换失败: returncode={result.returncode} stderr={result.stderr}")

        else:
            logger.info("未检测到 unoconv 可执行文件，尝试使用 libreoffice (soffice)")

        # 2) 尝试 libreoffice soffice（更普遍）
        soffice_path = shutil.which("soffice") or shutil.which("libreoffice")
        if soffice_path:
            # 使用 LibreOffice 命令行把文件转换到 output_dir，然后移动/重命名到 output_path
            outdir = output_dir if output_dir else os.path.dirname(os.path.abspath(output_path)) or "."
            cmd = [soffice_path, "--headless", "--convert-to", "pdf", "--outdir", outdir, input_path]
            logger.debug(f"调用 soffice: {' '.join(cmd)}")
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
            logger.debug(f"soffice returncode={result.returncode} stdout={result.stdout} stderr={result.stderr}")
            # soffice 会在 outdir 下生成与输入同名但扩展名为 .pdf 的文件
            src_pdf = os.path.join(outdir, os.path.splitext(os.path.basename(input_path))[0] + ".pdf")
            if result.returncode == 0 and os.path.exists(src_pdf):
                # 移动到期望的 output_path（覆盖）
                if os.path.abspath(src_pdf) != os.path.abspath(output_path):
                    try:
                        os.replace(src_pdf, output_path)
                    except Exception:
                        # 如果移动失败，尝试复制
                        import shutil as _shutil
                        _shutil.copyfile(src_pdf, output_path)
                logger.info(f"Fallback (soffice) PPT to PDF 成功: {output_path}")
                return True
            else:
                logger.warning(f"soffice 转换失败: returncode={result.returncode} stderr={result.stderr}")
        else:
            logger.error("未检测到 soffice/libreoffice 可执行文件，请安装 libreoffice 或 unoconv")

        logger.error("Fallback PPT to PDF 最终失败")
        return False

    except subprocess.TimeoutExpired:
        logger.error("Fallback PPT to PDF 超时")
        logger.debug(traceback.format_exc())
        return False
    except Exception as e:
        logger.error(f"Fallback PPT to PDF 出现错误: {str(e)}")
        logger.debug(traceback.format_exc())
        return False
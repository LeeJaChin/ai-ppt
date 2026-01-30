"""
PPT 生成引擎
使用 python-pptx 库生成格式化的 PowerPoint 文件
"""
import os
import logging
from typing import List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from ..models import SlideContent, ThemeStyle, SlideLayout


class PPTGenerator:
    """PPT 生成器类 - 极致视觉增强版"""
    
    # 主题配置
    THEMES = {
        ThemeStyle.BUSINESS: {
            "bg_color": RGBColor(255, 255, 255),
            "title_color": RGBColor(12, 45, 87),
            "text_color": RGBColor(40, 40, 40),
            "accent_color": RGBColor(255, 193, 7),
            "font_name": "微软雅黑",
            "decoration_shape": MSO_SHAPE.RECTANGLE
        },
        ThemeStyle.TECH: {
            "bg_color": RGBColor(10, 10, 26),
            "title_color": RGBColor(0, 255, 255),
            "text_color": RGBColor(200, 200, 220),
            "accent_color": RGBColor(138, 43, 226),
            "font_name": "微软雅黑",
            "decoration_shape": MSO_SHAPE.CHEVRON
        },
        ThemeStyle.CREATIVE: {
            "bg_color": RGBColor(255, 248, 240),
            "title_color": RGBColor(255, 87, 34),
            "text_color": RGBColor(62, 39, 35),
            "accent_color": RGBColor(76, 175, 80),
            "font_name": "微软雅黑",
            "decoration_shape": MSO_SHAPE.OVAL
        },
    }
    
    def __init__(self, theme: ThemeStyle = ThemeStyle.BUSINESS, template_path: str = None):
        self.logger = logging.getLogger("ai-ppt.generator")
        if template_path and os.path.exists(template_path):
            self.logger.info(f"正在从模板初始化 Presentation: {template_path}")
            try:
                self.prs = Presentation(template_path)
                self.template_mode = True
                # 记录模板中的布局名称，方便调试
                layout_names = [l.name for l in self.prs.slide_layouts]
                self.logger.debug(f"模板可用布局: {layout_names}")
            except Exception as e:
                self.logger.error(f"加载模板失败: {str(e)}，将使用默认样式")
                self.prs = Presentation()
                self.prs.slide_width = Inches(13.33)
                self.prs.slide_height = Inches(7.5)
                self.template_mode = False
        else:
            self.logger.info("正在初始化默认 Presentation")
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.33)
            self.prs.slide_height = Inches(7.5)
            self.template_mode = False
        
        self.theme = self.THEMES.get(theme, self.THEMES[ThemeStyle.BUSINESS])
    
    def _apply_font_style(self, run, size, color, bold=False):
        run.font.name = self.theme["font_name"]
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold

    def _setup_background(self, slide):
        if self.template_mode:
            return  # 模板模式下不强制覆盖背景
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme["bg_color"]

    def _add_page_header(self, slide, title, icon=None):
        """添加统一页眉"""
        if self.template_mode:
            # 模板模式下尝试查找已有的标题占位符
            for shape in slide.shapes:
                if not shape.is_placeholder:
                    continue
                if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    shape.text = f"{icon} {title}" if icon else title
                    return
            # 如果没通过类型找到，尝试通过名称兜底
            for shape in slide.shapes:
                if shape.is_placeholder and "TITLE" in str(shape.name).upper():
                    shape.text = f"{icon} {title}" if icon else title
                    return
            return

        # 装饰色块
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.5), Inches(0.4), Inches(0.2), Inches(0.8)
        )
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = self.theme["accent_color"]
        header_bar.line.fill.background()
        
        # 标题
        display_title = f"{icon} {title}" if icon else title
        title_shape = slide.shapes.add_textbox(Inches(0.9), Inches(0.35), Inches(12), Inches(1))
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = display_title
        self._apply_font_style(run, 36, self.theme["title_color"], bold=True)

    def add_title_slide(self, title: str, subtitle: str = "AI-PPT Architect 智绘大纲"):
        slide = self.prs.slides.add_slide(self._get_layout(is_title=True))
        self._setup_background(slide)
        
        # 装饰侧边 (仅非模板模式)
        if not self.template_mode:
            shape = slide.shapes.add_shape(
                self.theme["decoration_shape"], 
                Inches(0), Inches(0), Inches(0.8), self.prs.slide_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.theme["accent_color"]
            shape.line.fill.background()
        
        # 填充标题和副标题
        found_title = False
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue
            ph_type = shape.placeholder_format.type
            if ph_type == PP_PLACEHOLDER.TITLE or ph_type == PP_PLACEHOLDER.CENTER_TITLE:
                shape.text = title
                found_title = True
            elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                shape.text = subtitle
        
        if not found_title and not self.template_mode:
            # 手动添加 (仅在非模板模式或模板没标题位时兜底)
            txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(11), Inches(2.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
            self._apply_font_style(run, 60, self.theme["title_color"], bold=True)

    def _get_layout(self, is_title=False, layout_type: SlideLayout = None):
        """
        智能获取合适的布局
        优先匹配模板中的特定版式
        """
        if not self.template_mode:
            # 默认模式下的映射
            if is_title: return self.prs.slide_layouts[0]
            if layout_type == SlideLayout.TWO_COLUMN: return self.prs.slide_layouts[3] # Two Content
            return self.prs.slide_layouts[6] # Blank
        
        # 模板模式下寻找匹配的布局
        # 1. 根据 layout_type 尝试精确匹配
        search_keywords = []
        if is_title:
            search_keywords = ["TITLE SLIDE", "标题幻灯片", "封面", "TITLE"]
        elif layout_type == SlideLayout.TWO_COLUMN:
            search_keywords = ["TWO CONTENT", "两栏内容", "双栏", "COMPARISON", "对比"]
        elif layout_type in [SlideLayout.BULLETS, SlideLayout.PROCESS]:
            search_keywords = ["TITLE AND CONTENT", "标题和内容", "正文", "CONTENT"]
        elif layout_type == SlideLayout.THANK_YOU:
            search_keywords = ["THANK", "感谢", "结束", "CLOSING"]
        
        # 遍历布局名进行关键词匹配
        for keyword in search_keywords:
            for layout in self.prs.slide_layouts:
                if keyword in layout.name.upper():
                    self.logger.info(f"匹配到模板布局: {layout.name} (关键词: {keyword})")
                    return layout
        
        # 2. 如果没匹配到，尝试按索引匹配常用布局
        if is_title:
            return self.prs.slide_layouts[0]
        
        # 尝试寻找带 Content 占位符的布局
        for layout in self.prs.slide_layouts:
            for shape in layout.placeholders:
                if shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                    return layout
        
        # 兜底
        return self.prs.slide_layouts[1 if len(self.prs.slide_layouts) > 1 else 0]

    def add_bullet_slide(self, slide_data: SlideContent):
        slide = self.prs.slides.add_slide(self._get_layout(layout_type=SlideLayout.BULLETS))
        self._setup_background(slide)
        self._add_page_header(slide, slide_data.title, slide_data.icon)
        
        # 寻找正文占位符
        body_placeholder = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                body_placeholder = shape
                break
        
        if body_placeholder:
            tf = body_placeholder.text_frame
            tf.clear() # 清除默认文本
            for idx, point in enumerate(slide_data.bullet_points):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.text = point
                p.level = 0
        else:
            # 兜底手动添加
            left, top, width, height = Inches(1.2), Inches(1.8), Inches(11), Inches(4.5)
            body_shape = slide.shapes.add_textbox(left, top, width, height)
            tf = body_shape.text_frame
            tf.word_wrap = True
            
            for idx, point in enumerate(slide_data.bullet_points):
                p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                p.space_before = Pt(18)
                run = p.add_run()
                run.text = f"●  {point}"
                self._apply_font_style(run, 24, self.theme["text_color"])

    def add_column_slide(self, slide_data: SlideContent):
        slide = self.prs.slides.add_slide(self._get_layout(layout_type=SlideLayout.TWO_COLUMN))
        self._setup_background(slide)
        self._add_page_header(slide, slide_data.title, slide_data.icon)
        
        # 寻找双栏占位符
        placeholders = []
        for shape in sorted(slide.shapes, key=lambda s: s.left):
            if shape.is_placeholder and shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                placeholders.append(shape)
        
        # 拆分内容
        mid = len(slide_data.bullet_points) // 2
        left_points = slide_data.bullet_points[:mid]
        right_points = slide_data.bullet_points[mid:]
        
        if len(placeholders) >= 2:
            # 使用模板的双栏占位符
            for i, pts in enumerate([left_points, right_points]):
                tf = placeholders[i].text_frame
                tf.clear()
                for j, p_text in enumerate(pts):
                    p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
                    p.text = p_text
        else:
            # 左右双栏 手动添加
            col_width = Inches(5.5)
            mid_gap = Inches(0.5)
            left_start = Inches(1)
            
            for i, points in enumerate([left_points, right_points]):
                box = slide.shapes.add_textbox(
                    left_start + i*(col_width + mid_gap), 
                    Inches(2), col_width, Inches(4)
                )
                tf = box.text_frame
                tf.word_wrap = True
                for j, p_text in enumerate(points):
                    p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
                    p.space_before = Pt(15)
                    run = p.add_run()
                    run.text = f"▪ {p_text}"
                    self._apply_font_style(run, 20, self.theme["text_color"])

    def add_process_slide(self, slide_data: SlideContent):
        slide = self.prs.slides.add_slide(self._get_layout(layout_type=SlideLayout.PROCESS))
        self._setup_background(slide)
        self._add_page_header(slide, slide_data.title, slide_data.icon)
        
        # 流程图布局
        count = min(len(slide_data.bullet_points), 4)
        if count == 0: return
        
        shape_w = Inches(2.8)
        shape_h = Inches(1.5)
        start_x = (self.prs.slide_width - (shape_w * count + Inches(0.2) * (count-1))) / 2
        
        for i in range(count):
            x = start_x + i * (shape_w + Inches(0.2))
            # 步骤框
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3), shape_w, shape_h)
            rect.fill.solid()
            rect.fill.fore_color.rgb = self.theme["accent_color"]
            rect.line.color.rgb = self.theme["title_color"]
            
            # 步骤文字
            tf = rect.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = slide_data.bullet_points[i]
            self._apply_font_style(run, 18, self.theme["title_color"], bold=True)
            
            # 箭头
            if i < count - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, 
                    x + shape_w + Inches(0.02), Inches(3.5), Inches(0.16), Inches(0.5)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.theme["title_color"]

    def add_chart_slide(self, slide_data: SlideContent, chart_type: XL_CHART_TYPE):
        slide = self.prs.slides.add_slide(self._get_layout(layout_type=SlideLayout.DATA_COLUMN))
        self._setup_background(slide)
        self._add_page_header(slide, slide_data.title, slide_data.icon)
        
        # 布局参数
        has_text = len(slide_data.bullet_points) > 0
        
        # 尝试寻找占位符
        chart_placeholder = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type in [PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.BODY]:
                chart_placeholder = shape
                # 如果有文字，我们优先找右侧的占位符作为图表位
                if has_text and shape.left < self.prs.slide_width / 2:
                    continue
                break
        
        # 图表数据准备
        chart_data = ChartData()
        if slide_data.data_points:
            is_multi_series = any("series" in d for d in slide_data.data_points)
            if is_multi_series:
                labels = [d.get("label", f"项{i}") for i, d in enumerate(slide_data.data_points)]
                chart_data.categories = labels
                series_names = set()
                for d in slide_data.data_points:
                    if "series" in d:
                        series_names.update(d["series"].keys())
                for s_name in sorted(list(series_names)):
                    values = [d.get("series", {}).get(s_name, 0) for d in slide_data.data_points]
                    chart_data.add_series(s_name, values)
            else:
                chart_data.categories = [d.get("label", f"项{i}") for i, d in enumerate(slide_data.data_points)]
                chart_data.add_series("数值", [d.get("value", 0) for d in slide_data.data_points])
        else:
            chart_data.categories = ["示例 A", "示例 B", "示例 C"]
            chart_data.add_series("系列 1", [30, 50, 20])

        if chart_placeholder and not has_text:
            # 全屏图表占位符
            graphic_frame = chart_placeholder.insert_chart(chart_type, chart_data)
            chart = graphic_frame.chart
        else:
            # 采用手动布局或分栏布局
            if has_text:
                # 左侧文字区
                text_x, text_y, text_w, text_h = Inches(0.8), Inches(1.8), Inches(4.2), Inches(5)
                # 右侧图表区
                chart_x, chart_y, chart_w, chart_h = Inches(5.2), Inches(1.8), Inches(7.5), Inches(5)
                
                # 添加文字内容
                body_shape = slide.shapes.add_textbox(text_x, text_y, text_w, text_h)
                tf = body_shape.text_frame
                tf.word_wrap = True
                for idx, point in enumerate(slide_data.bullet_points):
                    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                    p.space_before = Pt(12)
                    run = p.add_run()
                    run.text = f"• {point}"
                    self._apply_font_style(run, 18, self.theme["text_color"])
            else:
                # 全屏图表 手动
                chart_x, chart_y, chart_w, chart_h = Inches(1.5), Inches(1.8), Inches(10.5), Inches(5)

            graphic_frame = slide.shapes.add_chart(
                chart_type, chart_x, chart_y, chart_w, chart_h, chart_data
            )
            chart = graphic_frame.chart
        
        # 图表样式微调
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        
        # 字体应用 (仅当非模板模式或手动添加时)
        if not self.template_mode and chart.has_title:
            self._apply_font_style(chart.chart_title.text_frame.paragraphs[0].runs[0], 14, self.theme["title_color"])

    def add_timeline_slide(self, slide_data: SlideContent):
        """添加时间轴/里程碑页"""
        slide = self.prs.slides.add_slide(self._get_layout(layout_type=SlideLayout.TIMELINE))
        self._setup_background(slide)
        self._add_page_header(slide, slide_data.title, slide_data.icon)
        
        # 时间轴主线
        line_y = Inches(4)
        line_start = Inches(1)
        line_end = self.prs.slide_width - Inches(1)
        
        main_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, line_start, line_y, line_end - line_start, Pt(4)
        )
        main_line.fill.solid()
        main_line.fill.fore_color.rgb = self.theme["accent_color"]
        main_line.line.fill.background()
        
        # 节点数量
        points = slide_data.bullet_points or ["开始", "过程", "结束"]
        count = len(points)
        if count == 0: return
        
        gap = (line_end - line_start) / max(count - 1, 1)
        
        for i, point in enumerate(points):
            x = line_start + i * gap
            
            # 节点圆圈
            r = Inches(0.2)
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - r/2, line_y - r/2 + Pt(2), r, r)
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.theme["title_color"]
            circle.line.color.rgb = self.theme["accent_color"]
            circle.line.width = Pt(2)
            
            # 节点文字 (交错上下显示)
            is_top = i % 2 == 0
            text_y = line_y - Inches(1.2) if is_top else line_y + Inches(0.5)
            
            box = slide.shapes.add_textbox(x - Inches(1), text_y, Inches(2), Inches(0.8))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = point
            self._apply_font_style(run, 16, self.theme["text_color"], bold=True)

    def add_big_number_slide(self, slide_data: SlideContent):
        """添加数字大屏/关键KPI页"""
        slide = self.prs.slides.add_slide(self._get_layout(layout_type=SlideLayout.BIG_NUMBER))
        self._setup_background(slide)
        self._add_page_header(slide, slide_data.title, slide_data.icon)
        
        center_x, center_y = self.prs.slide_width / 2, self.prs.slide_height / 2
        
        # 提取第一个数字
        big_val = "100%"
        if slide_data.data_points:
            first_dp = slide_data.data_points[0]
            big_val = str(first_dp.get("value", first_dp.get("label", "100%")))
            
        # 渲染大数字
        txBox = slide.shapes.add_textbox(center_x - Inches(3), center_y - Inches(1.5), Inches(6), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = big_val
        self._apply_font_style(run, 120, self.theme["accent_color"], bold=True)
        
        # 下方说明文字
        if slide_data.bullet_points:
            descBox = slide.shapes.add_textbox(center_x - Inches(4), center_y + Inches(1), Inches(8), Inches(1.5))
            tf_desc = descBox.text_frame
            tf_desc.word_wrap = True
            p_desc = tf_desc.paragraphs[0]
            p_desc.alignment = PP_ALIGN.CENTER
            run_desc = p_desc.add_run()
            run_desc.text = slide_data.bullet_points[0]
            self._apply_font_style(run_desc, 24, self.theme["text_color"])

    def add_thank_you_slide(self, message: str = "感谢聆听"):
        """添加精美致谢页"""
        slide = self.prs.slides.add_slide(self._get_layout(layout_type=SlideLayout.THANK_YOU))
        self._setup_background(slide)
        
        # 尝试填充占位符
        found_placeholder = False
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.BODY]:
                shape.text = message
                found_placeholder = True
                break
        
        if found_placeholder:
            return

        center_x, center_y = self.prs.slide_width / 2, self.prs.slide_height / 2
        
        # 装饰大框
        box_w, box_h = Inches(8), Inches(3)
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            center_x - box_w/2, center_y - box_h/2, box_w, box_h
        )
        box.fill.background()
        box.line.color.rgb = self.theme["accent_color"]
        box.line.width = Pt(5)
        
        txBox = slide.shapes.add_textbox(center_x - box_w/2, center_y - Inches(0.6), box_w, Inches(1.2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = message
        self._apply_font_style(run, 64, self.theme["title_color"], bold=True)

    def generate(self, title: str, slides: List[SlideContent], output_path: str) -> str:
        try:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            
            # 如果是模板模式，且模板本身不是空的（即有超过0页），我们通常是在后面追加。
            # 但用户通常希望“基于模板”生成，如果模板只有母版而没有页面，则直接开始。
            # 如果模板有封面页，我们甚至可以考虑直接修改封面页。
            
            has_existing_slides = len(self.prs.slides) > 0
            
            if self.template_mode and has_existing_slides:
                # 尝试寻找并填充已有的封面
                first_slide = self.prs.slides[0]
                found_title = False
                for shape in first_slide.shapes:
                    if shape.is_placeholder and shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                        shape.text = title
                        found_title = True
                
                # 如果没在第一页找到标题位，且模板模式开启，通常我们不主动增加新封面，以免破坏模板结构
                # 除非用户明确要求（目前逻辑是跳过 TITLE 布局的循环）
            elif not self.template_mode:
                self.add_title_slide(title)
            else:
                # 模板模式但没页面，还是得加个封面
                self.add_title_slide(title)
            
            for slide in slides:
                if slide.layout == SlideLayout.TITLE:
                    continue
                elif slide.layout == SlideLayout.TWO_COLUMN:
                    self.add_column_slide(slide)
                elif slide.layout == SlideLayout.PROCESS:
                    self.add_process_slide(slide)
                elif slide.layout == SlideLayout.DATA_COLUMN:
                    self.add_chart_slide(slide, XL_CHART_TYPE.COLUMN_CLUSTERED)
                elif slide.layout == SlideLayout.DATA_BAR:
                    self.add_chart_slide(slide, XL_CHART_TYPE.BAR_CLUSTERED)
                elif slide.layout == SlideLayout.DATA_LINE:
                    self.add_chart_slide(slide, XL_CHART_TYPE.LINE)
                elif slide.layout == SlideLayout.DATA_PIE:
                    self.add_chart_slide(slide, XL_CHART_TYPE.PIE)
                elif slide.layout == SlideLayout.DATA_AREA:
                    self.add_chart_slide(slide, XL_CHART_TYPE.AREA)
                elif slide.layout == SlideLayout.DATA_STACKED:
                    self.add_chart_slide(slide, XL_CHART_TYPE.COLUMN_STACKED)
                elif slide.layout == SlideLayout.TIMELINE:
                    self.add_timeline_slide(slide)
                elif slide.layout == SlideLayout.BIG_NUMBER:
                    self.add_big_number_slide(slide)
                elif slide.layout == SlideLayout.THANK_YOU:
                    self.add_thank_you_slide(slide.title)
                else:
                    self.add_bullet_slide(slide)
            self.prs.save(output_path)
            return output_path
        except Exception as e:
            raise Exception(f"PPT 生成失败: {str(e)}")

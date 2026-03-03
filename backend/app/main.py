"""
FastAPI 主应用入口 - 已修改去除水印 (使用 LibreOffice)
"""
import os
import uuid
import asyncio
import logging
import subprocess  # 新增: 用于调用外部命令
import sys        # 新增: 用于判断操作系统
import shutil     # 新增: 用于移动文件
from datetime import datetime
from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from typing import Dict

from pdf2docx import Converter

from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt
import io

# 保持原有的 import，但注意 converter 可能不再被完全依赖，除非用来做其他格式转换
from .config import settings
from .models import (
    GenerateOutlineRequest,
    GeneratePPTRequest,
    OutlineResponse,
    TaskResponse,
    TaskStatus,
    UserCreate,
    UserLogin,
    User,
    Token,
    HistoryResponse,
)
from .services.ai_factory import AIAdapterFactory
from .services.ppt_generator import PPTGenerator
from .services.converter import FileConverter
from .services.redis_client import redis_client
from .services.auth import AuthService

# 配置日志
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("ai-ppt")

# 创建 FastAPI 应用
app = FastAPI(
    title="AI-PPT Architect API",
    description="AI 驱动的 PPT 自动生成服务",
    version="1.0.0",
)

# 配置 CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=settings.cors_origins_list,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 任务存储
tasks_storage: Dict[str, Dict] = {}

# --- 新增函数: 使用 LibreOffice 进行无水印转换 ---
# --- 替换 main.py 中的 convert_with_libreoffice 函数 ---

def convert_with_libreoffice(input_path: str, output_path: str) -> bool:
    """
    使用 LibreOffice 将 PPT/Word 转换为 PDF (无水印)
    兼容 Windows, macOS, Linux
    """
    try:
        # 1. 获取绝对路径 (防止相对路径在子进程中出错)
        input_abs_path = os.path.abspath(input_path)
        output_dir_abs = os.path.abspath(os.path.dirname(output_path))

        # 2. 确定 LibreOffice 的执行命令路径
        soffice_cmd = "libreoffice" # Linux 默认

        if sys.platform == "darwin": # macOS
            # macOS 标准安装路径
            mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
            if os.path.exists(mac_path):
                soffice_cmd = mac_path
            else:
                # 尝试查找用户可能手动配置在 PATH 中的命令
                if shutil.which("libreoffice"):
                    soffice_cmd = "libreoffice"
                elif shutil.which("soffice"):
                    soffice_cmd = "soffice"
                else:
                    logger.error("未找到 LibreOffice。请确保已安装: brew install --cask libreoffice")
                    return False

        elif sys.platform == "win32": # Windows
            possible_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
            ]
            # 先检查默认路径
            found = False
            for p in possible_paths:
                if os.path.exists(p):
                    soffice_cmd = p
                    found = True
                    break
            # 如果默认路径没有，尝试环境变量
            if not found:
                soffice_cmd = "soffice"

        # 3. 构建命令
        # 注意：使用绝对路径
        cmd = [
            soffice_cmd,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_dir_abs,
            input_abs_path
        ]

        logger.info(f"执行转换命令: {cmd}")

        # 4. 执行转换
        result = subprocess.run(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=300 # 增加超时时间，PPT转PDF可能较慢
        )

        if result.returncode != 0:
            logger.error(f"LibreOffice 转换失败 (Code {result.returncode}): {result.stderr.decode()}")
            return False

        # 5. 处理文件名重命名
        # LibreOffice 默认输出文件名是 [原文件名].pdf
        input_filename_no_ext = os.path.splitext(os.path.basename(input_path))[0]
        # 注意：LibreOffice 可能会根据原文件自动去除 .pptx 后缀生成 .pdf
        expected_output_name = f"{input_filename_no_ext}.pdf"
        generated_file_path = os.path.join(output_dir_abs, expected_output_name)

        if os.path.exists(generated_file_path):
            # 只有当生成的文件名和目标文件名不一样时才重命名
            # 注意：os.path.samefile 在文件不存在时会报错，所以直接比较路径字符串
            if os.path.abspath(output_path) != generated_file_path:
                # 如果目标文件已存在，先删除，防止 windows 下报错
                if os.path.exists(output_path):
                    os.remove(output_path)
                shutil.move(generated_file_path, output_path)

            logger.info(f"转换成功: {output_path}")
            return True
        else:
            logger.error(f"未找到生成的 PDF 文件，预期路径: {generated_file_path}")
            # 调试：列出目录下文件
            logger.error(f"目录下现有文件: {os.listdir(output_dir_abs)}")
            return False

    except Exception as e:
        logger.error(f"LibreOffice 转换异常: {str(e)}", exc_info=True)
        return False
# ----------------------------------------------------


def convert_pdf_to_pptx_file(input_path: str, output_path: str) -> bool:
    """
    将 PDF 转换为 PPTX (通过将每一页转为图片的方式)
    """
    try:
        logger.info(f"开始 PDF 转 PPT: {input_path}")

        # 1. 将 PDF 转为图片列表
        # thread_count 指定线程数，提高速度
        images = convert_from_path(input_path, thread_count=2)

        if not images:
            logger.error("未从 PDF 解析出任何页面")
            return False

        # 2. 创建一个新的 PPT 对象
        prs = Presentation()

        # 删除默认的第一张空白幻灯片(如果有)
        if len(prs.slides) > 0:
            # python-pptx 不直接支持删除 slide，但新建的默认通常是空的或者带一个标题页
            # 我们不需要管它，直接设置 slide layout 即可，或者用空白模版
            pass

        # 获取第一张图片的尺寸，用于设置 PPT 幻灯片大小
        # 图片是 PIL Image 对象
        width, height = images[0].size

        # 设置 PPT 的页面大小与 PDF 页面比例一致
        # python-pptx 默认单位是 EMU (English Metric Unit)，我们需要转换
        # 这里为了简单，我们根据图片像素比例设置幻灯片尺寸
        # 默认 PPT 宽 10 英寸，高 7.5 英寸

        # 简单的做法：固定 PPT 宽度为 10 英寸，高度按比例缩放
        aspect_ratio = height / width
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(10 * aspect_ratio)

        # 3. 遍历图片，添加到幻灯片
        blank_slide_layout = prs.slide_layouts[6] # 6 是空白布局

        for i, img in enumerate(images):
            # 添加一张新幻灯片
            slide = prs.slides.add_slide(blank_slide_layout)

            # 将 PIL image 转为字节流
            image_stream = io.BytesIO()
            img.save(image_stream, format='PNG')
            image_stream.seek(0)

            # 将图片铺满整个幻灯片
            slide.shapes.add_picture(
                image_stream,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )
            logger.info(f"处理第 {i+1} 页...")

        # 4. 保存 PPT
        prs.save(output_path)
        logger.info(f"PPT 生成成功: {output_path}")
        return True

    except Exception as e:
        logger.error(f"PDF 转 PPT 失败: {str(e)}", exc_info=True)
        return False

@app.get("/")
async def root():
    return {"message": "AI-PPT Architect API", "status": "running", "version": "1.0.0"}

@app.get("/api/models")
async def get_available_models():
    try:
        models = AIAdapterFactory.get_available_models()
        return {"models": models, "count": len(models)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/upload-template")
async def upload_template(file: UploadFile = File(...)):
    if not file.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="只支持 .pptx 格式的模板文件")
    try:
        os.makedirs(settings.templates_dir, exist_ok=True)
        template_id = str(uuid.uuid4())
        file_ext = os.path.splitext(file.filename)[1]
        filename = f"{template_id}{file_ext}"
        save_path = os.path.join(settings.templates_dir, filename)
        with open(save_path, "wb") as f:
            content = await file.read()
            f.write(content)
        logger.info(f"模板上传成功: {file.filename} -> {template_id}")
        return {"template_id": template_id, "filename": file.filename, "message": "模板上传成功"}
    except Exception as e:
        logger.error(f"模板上传失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"模板上传失败: {str(e)}")

@app.post("/api/generate-outline", response_model=OutlineResponse)
async def generate_outline(request: GenerateOutlineRequest):
    logger.info(f"收到大纲生成请求: 模型={request.model.value}, 内容长度={len(request.content)}")
    try:
        adapter = AIAdapterFactory.create_adapter(request.model)
        outline_data = await adapter.generate_outline(request.content, slide_count=request.slide_count)
        return OutlineResponse(**outline_data)
    except ValueError as e:
        logger.error(f"参数错误: {str(e)}")
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error(f"大纲生成异常: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"大纲生成失败: {str(e)}")

@app.post("/api/generate-ppt", response_model=TaskResponse)
async def generate_ppt(request: GeneratePPTRequest):
    try:
        task_id = str(uuid.uuid4())
        task_data = {
            "status": TaskStatus.PENDING,
            "progress": 0,
            "message": "任务已创建",
            "created_at": datetime.now().isoformat(),
        }
        # 存储到内存
        tasks_storage[task_id] = task_data
        # 存储到Redis
        redis_client.set(f"task:{task_id}", task_data)
        asyncio.create_task(process_ppt_generation(task_id, request))
        return TaskResponse(
            task_id=task_id,
            status=TaskStatus.PENDING,
            progress=0,
            message="PPT 生成任务已启动"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"任务创建失败: {str(e)}")

import concurrent.futures

async def process_ppt_generation(task_id: str, request: GeneratePPTRequest):
    try:
        # 更新任务状态
        def update_task_status(status, progress, message, **kwargs):
            tasks_storage[task_id]["status"] = status
            tasks_storage[task_id]["progress"] = progress
            tasks_storage[task_id]["message"] = message
            for key, value in kwargs.items():
                tasks_storage[task_id][key] = value
            # 同时更新Redis
            redis_client.set(f"task:{task_id}", tasks_storage[task_id])

        update_task_status(TaskStatus.PROCESSING, 10, "正在初始化...")
        os.makedirs(settings.output_dir, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"ppt_{timestamp}_{task_id[:8]}.pptx"
        output_path = os.path.join(settings.output_dir, filename)

        update_task_status(TaskStatus.PROCESSING, 30, "正在生成 PPT...")

        template_path = None
        if request.template_id:
            template_path = os.path.abspath(os.path.join(settings.templates_dir, f"{request.template_id}.pptx"))
            if not os.path.exists(template_path):
                template_path = None

        # 使用线程池执行器来避免阻塞事件循环
        loop = asyncio.get_event_loop()
        with concurrent.futures.ThreadPoolExecutor() as executor:
            def generate_ppt_sync():
                generator = PPTGenerator(theme=request.theme, template_path=template_path)
                return generator.generate(
                    title=request.outline.title,
                    slides=request.outline.slides,
                    output_path=output_path
                )
            
            file_path = await loop.run_in_executor(executor, generate_ppt_sync)

        update_task_status(TaskStatus.PROCESSING, 90, "正在完成...")
        await asyncio.sleep(0.5)
        update_task_status(
            TaskStatus.COMPLETED, 
            100, 
            "PPT 生成完成",
            file_path=file_path,
            download_url=f"/api/download/{task_id}"
        )
        
        # 尝试从请求中获取用户信息并添加历史记录
        # 注意：这里简化处理，实际应该从请求上下文中获取用户信息
        # 由于当前实现中没有传递用户信息，这里暂时注释掉
        # AuthService.add_history_record(user_id, request.outline.title, task_id, file_path)
        
    except Exception as e:
        logger.error(f"PPT生成异常: {str(e)}", exc_info=True)
        error_data = {
            "status": TaskStatus.FAILED,
            "message": f"生成失败: {str(e)}",
            "error": str(e)
        }
        tasks_storage[task_id].update(error_data)
        redis_client.set(f"task:{task_id}", tasks_storage[task_id])

@app.get("/api/task/{task_id}", response_model=TaskResponse)
async def get_task_status(task_id: str):
    # 首先从内存中查找
    if task_id in tasks_storage:
        task = tasks_storage[task_id]
    else:
        # 如果内存中没有，从Redis中查找
        task = redis_client.get(f"task:{task_id}")
        if not task:
            raise HTTPException(status_code=404, detail="任务不存在")
        # 将Redis中的任务数据加载到内存中
        tasks_storage[task_id] = task
    return TaskResponse(
        task_id=task_id,
        status=task["status"],
        progress=task["progress"],
        message=task.get("message"),
        download_url=task.get("download_url")
    )

@app.get("/api/download/{task_id}")
async def download_ppt(task_id: str):
    # 首先从内存中查找
    if task_id in tasks_storage:
        task = tasks_storage[task_id]
    else:
        # 如果内存中没有，从Redis中查找
        task = redis_client.get(f"task:{task_id}")
        if not task:
            raise HTTPException(status_code=404, detail="任务不存在")
    if task["status"] != TaskStatus.COMPLETED:
        raise HTTPException(status_code=400, detail="文件尚未生成完成")
    file_path = task.get("file_path")
    if not file_path or not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="文件不存在")
    filename = os.path.basename(file_path)
    return FileResponse(
        path=file_path,
        filename=filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

@app.get("/api/convert/ppt-to-pdf/{task_id}")
async def convert_ppt_to_pdf(task_id: str):
    """将生成的PPT转换为PDF用于在线预览"""
    try:
        # 首先从内存中查找
        if task_id in tasks_storage:
            task = tasks_storage[task_id]
        else:
            # 如果内存中没有，从Redis中查找
            task = redis_client.get(f"task:{task_id}")
            if not task:
                raise HTTPException(status_code=404, detail="任务不存在")
        
        if task["status"] != TaskStatus.COMPLETED:
            raise HTTPException(status_code=400, detail="PPT尚未生成完成")
        
        ppt_path = task.get("file_path")
        if not ppt_path or not os.path.exists(ppt_path):
            raise HTTPException(status_code=404, detail="PPT文件不存在")
        
        # 生成PDF路径
        pdf_path = os.path.splitext(ppt_path)[0] + ".pdf"
        
        # 检查PDF是否已经存在
        if os.path.exists(pdf_path):
            # 如果存在，直接返回
            task_id = str(uuid.uuid4())
            tasks_storage[task_id] = {
                "status": TaskStatus.COMPLETED,
                "progress": 100,
                "message": "PDF预览已生成",
                "file_path": pdf_path,
                "download_url": f"/api/download/{task_id}"
            }
            redis_client.set(f"task:{task_id}", tasks_storage[task_id])
            return {
                "task_id": task_id,
                "status": TaskStatus.COMPLETED,
                "progress": 100,
                "message": "PDF预览已生成",
                "download_url": f"/api/download/{task_id}"
            }
        
        # 调用LibreOffice转换
        success = convert_with_libreoffice(ppt_path, pdf_path)
        
        if success:
            task_id = str(uuid.uuid4())
            tasks_storage[task_id] = {
                "status": TaskStatus.COMPLETED,
                "progress": 100,
                "message": "PDF预览生成完成",
                "file_path": pdf_path,
                "download_url": f"/api/download/{task_id}"
            }
            redis_client.set(f"task:{task_id}", tasks_storage[task_id])
            return {
                "task_id": task_id,
                "status": TaskStatus.COMPLETED,
                "progress": 100,
                "message": "PDF预览生成完成",
                "download_url": f"/api/download/{task_id}"
            }
        else:
            raise HTTPException(status_code=500, detail="转换PDF失败")
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"PPT转PDF失败: {str(e)}")
        raise HTTPException(status_code=500, detail=f"转换PDF失败: {str(e)}")

@app.post("/api/convert", response_model=TaskResponse)
async def convert_file(file: UploadFile = File(...), target_format: str = "pdf"):
    ext = os.path.splitext(file.filename)[1].lower()
    task_id = str(uuid.uuid4())
    task_data = {
        "status": TaskStatus.PENDING,
        "progress": 0,
        "message": "转换任务已启动",
        "created_at": datetime.now().isoformat(),
    }
    # 存储到内存
    tasks_storage[task_id] = task_data
    # 存储到Redis
    redis_client.set(f"task:{task_id}", task_data)
    temp_dir = os.path.join(settings.output_dir, "temp")
    os.makedirs(temp_dir, exist_ok=True)
    input_filename = f"{task_id}_in{ext}"
    input_path = os.path.join(temp_dir, input_filename)
    with open(input_path, "wb") as f:
        content = await file.read()
        f.write(content)
    output_ext = f".{target_format}" if not target_format.startswith(".") else target_format
    output_filename = f"{task_id}_out{output_ext}"
    output_path = os.path.join(settings.output_dir, output_filename)
    asyncio.create_task(process_conversion(task_id, input_path, output_path, ext, output_ext))
    return TaskResponse(task_id=task_id, status=TaskStatus.PENDING, progress=0, message="文件转换中...")

# --- 修改后的转换任务处理函数 ---
# --- 修改后的 process_conversion 函数 ---
# --- 更新后的 process_conversion ---
async def process_conversion(task_id: str, input_path: str, output_path: str, in_ext: str, out_ext: str):
    """处理文件转换后台任务"""
    try:
        # 更新任务状态的函数
        def update_task_status(status, progress, message, **kwargs):
            tasks_storage[task_id]["status"] = status
            tasks_storage[task_id]["progress"] = progress
            tasks_storage[task_id]["message"] = message
            for key, value in kwargs.items():
                tasks_storage[task_id][key] = value
            # 同时更新Redis
            redis_client.set(f"task:{task_id}", tasks_storage[task_id])

        update_task_status(TaskStatus.PROCESSING, 20, "正在处理文件...")

        success = False

        # 1. PPT/Word -> PDF (使用 LibreOffice)
        if out_ext == '.pdf' and in_ext in ['.ppt', '.pptx', '.doc', '.docx']:
            update_task_status(TaskStatus.PROCESSING, 40, "正在转换为PDF...")
            success = convert_with_libreoffice(input_path, output_path)

        # 2. PDF -> Word (使用 pdf2docx)
        elif in_ext == '.pdf' and out_ext in ['.docx', '.doc']:
            update_task_status(TaskStatus.PROCESSING, 40, "正在转换为Word...")
            real_output = output_path
            if output_path.endswith('.doc'):
                real_output = output_path + 'x'
            success = convert_pdf_to_docx_file(input_path, real_output)
            if success and real_output != output_path:
                shutil.move(real_output, output_path)

        # 3. PDF -> PPT (使用 pdf2image + python-pptx) <--- 只有这里变了
        elif in_ext == '.pdf' and out_ext in ['.pptx', '.ppt']:
            update_task_status(TaskStatus.PROCESSING, 40, "正在转换为PPT...")
            # 处理 .ppt 后缀兼容
            real_output = output_path
            if output_path.endswith('.ppt'):
                real_output = output_path + 'x'

            success = convert_pdf_to_pptx_file(input_path, real_output)

            if success and real_output != output_path:
                shutil.move(real_output, output_path)

        else:
            logger.error(f"不支持的转换类型: {in_ext} -> {out_ext}")
            raise HTTPException(status_code=400, detail=f"不支持的转换类型: {in_ext} to {out_ext}")

        if success:
            update_task_status(
                TaskStatus.COMPLETED, 
                100, 
                "转换完成",
                file_path=output_path,
                download_url=f"/api/download/{task_id}"
            )
        else:
            raise Exception("转换未能生成目标文件")

    except Exception as e:
        error_data = {
            "status": TaskStatus.FAILED,
            "message": f"转换失败: {str(e)}"
        }
        tasks_storage[task_id].update(error_data)
        redis_client.set(f"task:{task_id}", tasks_storage[task_id])
        logger.error(f"转换任务 {task_id} 失败: {str(e)}", exc_info=True)
    finally:
        # 清理临时输入文件
        if os.path.exists(input_path):
            try:
                os.remove(input_path)
            except:
                pass

def convert_pdf_to_docx_file(input_path: str, output_path: str) -> bool:
    """
    使用 pdf2docx 库将 PDF 转换为 Word
    """
    try:
        logger.info(f"开始 PDF 转 Word: {input_path}")

        # 使用 pdf2docx 进行转换
        cv = Converter(input_path)
        # start=0, end=None 表示转换所有页面
        cv.convert(output_path, start=0, end=None)
        cv.close()

        if os.path.exists(output_path):
            logger.info(f"PDF 转 Word 成功: {output_path}")
            return True
        else:
            logger.error("文件未生成")
            return False

    except Exception as e:
        logger.error(f"PDF 转 Word 失败: {str(e)}", exc_info=True)
        return False

# 认证相关的API端点
@app.post("/api/auth/register")
async def register(user_data: UserCreate):
    """用户注册"""
    try:
        # 检查用户是否已存在
        existing_user = AuthService.get_user_by_email(user_data.email)
        if existing_user:
            raise HTTPException(status_code=400, detail="该邮箱已被注册")
        
        # 创建新用户
        user = AuthService.create_user(user_data.name, user_data.email, user_data.password)
        
        # 生成访问令牌
        access_token = AuthService.create_access_token(data={"sub": user["id"]})
        
        return Token(
            access_token=access_token,
            token_type="bearer",
            user=User(
                id=user["id"],
                name=user["name"],
                email=user["email"],
                is_active=user["is_active"],
                created_at=user["created_at"]
            )
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"注册失败: {str(e)}")
        raise HTTPException(status_code=500, detail="注册失败")

@app.post("/api/auth/login")
async def login(login_data: UserLogin):
    """用户登录"""
    try:
        # 查找用户
        user = AuthService.get_user_by_email(login_data.email)
        if not user:
            raise HTTPException(status_code=401, detail="邮箱或密码错误")
        
        # 验证密码
        if not AuthService.verify_password(login_data.password, user["password"]):
            raise HTTPException(status_code=401, detail="邮箱或密码错误")
        
        # 生成访问令牌
        access_token = AuthService.create_access_token(data={"sub": user["id"]})
        
        return Token(
            access_token=access_token,
            token_type="bearer",
            user=User(
                id=user["id"],
                name=user["name"],
                email=user["email"],
                is_active=user["is_active"],
                created_at=user["created_at"]
            )
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"登录失败: {str(e)}")
        raise HTTPException(status_code=500, detail="登录失败")

@app.get("/api/auth/me")
async def get_current_user(request):
    """获取当前用户信息"""
    try:
        # 从请求头获取令牌
        auth_header = request.headers.get("Authorization")
        
        if not auth_header:
            raise HTTPException(status_code=401, detail="未提供认证令牌")
        
        token = auth_header.replace("Bearer ", "")
        payload = AuthService.decode_token(token)
        
        if not payload:
            raise HTTPException(status_code=401, detail="无效的认证令牌")
        
        user_id = payload.get("sub")
        if not user_id:
            raise HTTPException(status_code=401, detail="无效的认证令牌")
        
        # 从auth模块导入users_db
        from app.services.auth import users_db
        
        # 查找用户
        user = users_db.get(user_id)
        if not user:
            raise HTTPException(status_code=401, detail="用户不存在")
        
        return User(
            id=user["id"],
            name=user["name"],
            email=user["email"],
            is_active=user["is_active"],
            created_at=user["created_at"]
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"获取用户信息失败: {str(e)}")
        raise HTTPException(status_code=500, detail="获取用户信息失败")

@app.get("/api/history")
async def get_history(request):
    """获取用户历史记录"""
    try:
        # 从请求头获取令牌
        auth_header = request.headers.get("Authorization")
        
        if not auth_header:
            raise HTTPException(status_code=401, detail="未提供认证令牌")
        
        token = auth_header.replace("Bearer ", "")
        payload = AuthService.decode_token(token)
        
        if not payload:
            raise HTTPException(status_code=401, detail="无效的认证令牌")
        
        user_id = payload.get("sub")
        if not user_id:
            raise HTTPException(status_code=401, detail="无效的认证令牌")
        
        # 获取用户历史记录
        history = AuthService.get_user_history(user_id)
        
        return HistoryResponse(
            items=[
                HistoryItem(
                    id=item["id"],
                    user_id=item["user_id"],
                    title=item["title"],
                    task_id=item["task_id"],
                    file_path=item.get("file_path"),
                    created_at=item["created_at"]
                ) for item in history
            ],
            total=len(history)
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"获取历史记录失败: {str(e)}")
        raise HTTPException(status_code=500, detail="获取历史记录失败")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(
        "app.main:app",
        host=settings.api_host,
        port=settings.api_port,
        reload=True
    )